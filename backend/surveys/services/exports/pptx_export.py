from io import BytesIO

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


def _rgb_from_hex(value):
    hex_value = (value or "#2563eb").strip().lstrip("#")
    if len(hex_value) == 3:
        hex_value = "".join(character * 2 for character in hex_value)
    hex_value = (hex_value[:6] or "2563EB").ljust(6, "0")
    return RGBColor(
        int(hex_value[0:2], 16),
        int(hex_value[2:4], 16),
        int(hex_value[4:6], 16),
    )


class PPTXExportService:
    def generate(self, survey, analytics_data, config) -> bytes:
        presentation = Presentation()
        presentation.slide_width = Inches(13.333)
        presentation.slide_height = Inches(7.5)

        self.brand_color = _rgb_from_hex(analytics_data["branding"]["color"])
        self._add_title_slide(presentation, survey, analytics_data)
        self._add_summary_slide(presentation, analytics_data)

        for section in analytics_data["questions"]:
            self._add_question_slide(presentation, section)

        for cross_tab in analytics_data["cross_tabs"]:
            self._add_cross_tab_slide(presentation, cross_tab)

        self._add_closing_slide(presentation, analytics_data)

        output = BytesIO()
        presentation.save(output)
        return output.getvalue()

    def _add_title_slide(self, presentation, survey, analytics_data):
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        self._add_footer_bar(slide)
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.4), Inches(9), Inches(1))
        title_frame = title_box.text_frame
        title_frame.text = survey.title
        title_frame.paragraphs[0].font.size = Pt(28)
        title_frame.paragraphs[0].font.bold = True

        subtitle_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.4), Inches(9), Inches(1.2))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = (
            f"{analytics_data['summary']['total_responses']} responses | "
            f"Generated {analytics_data['generated_at'].strftime('%Y-%m-%d %H:%M UTC')}"
        )
        subtitle_frame.paragraphs[0].font.size = Pt(16)

    def _add_summary_slide(self, presentation, analytics_data):
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        self._add_footer_bar(slide)
        self._add_slide_title(slide, "Executive Summary")
        summary = analytics_data["summary"]

        metric_specs = [
            ("Total Responses", summary["total_responses"]),
            ("Completion Rate", f"{summary['completion_rate']}%"),
            ("Average Duration", f"{summary['average_duration_seconds']}s"),
            ("In Progress", summary["in_progress_count"]),
        ]
        positions = [
            (Inches(0.8), Inches(1.7)),
            (Inches(6.9), Inches(1.7)),
            (Inches(0.8), Inches(4.0)),
            (Inches(6.9), Inches(4.0)),
        ]
        for (label, value), (left, top) in zip(metric_specs, positions):
            self._add_metric_box(slide, left, top, Inches(5.3), Inches(1.8), label, value)

    def _add_question_slide(self, presentation, section):
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        self._add_footer_bar(slide)
        self._add_slide_title(slide, section["text"])
        self._add_stats_box(slide, section)

        if section["analytics_type"] in {"categorical", "numeric", "temporal"} and section["table_rows"]:
            self._add_chart(slide, section)
        else:
            self._add_table_or_text(slide, section)

    def _add_cross_tab_slide(self, presentation, cross_tab):
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        self._add_footer_bar(slide)
        self._add_slide_title(slide, f"Cross-tabulation: {cross_tab['title']}")
        rows = len(cross_tab["matrix"]) + 2
        cols = len(cross_tab["col_totals"]) + 2
        table = slide.shapes.add_table(rows, cols, Inches(0.8), Inches(1.6), Inches(11.6), Inches(4.8)).table

        table.cell(0, 0).text = "Row label"
        for index, total in enumerate(cross_tab["col_totals"], start=1):
            table.cell(0, index).text = total["col_label"]
        table.cell(0, cols - 1).text = "Total"

        for row_index, row in enumerate(cross_tab["matrix"], start=1):
            table.cell(row_index, 0).text = row["row_label"]
            for col_index, cell in enumerate(row["cells"], start=1):
                table.cell(row_index, col_index).text = str(cell["count"])
            table.cell(row_index, cols - 1).text = str(row["row_total"])

        table.cell(rows - 1, 0).text = "Column total"
        for index, total in enumerate(cross_tab["col_totals"], start=1):
            table.cell(rows - 1, index).text = str(total["total"])
        table.cell(rows - 1, cols - 1).text = str(cross_tab["grand_total"])

        chi_box = slide.shapes.add_textbox(Inches(0.8), Inches(6.55), Inches(8), Inches(0.5))
        chi_frame = chi_box.text_frame
        chi_frame.text = (
            f"Chi-square: {cross_tab['chi_square']['statistic']} | "
            f"p-value: {cross_tab['chi_square']['p_value']}"
        )
        chi_frame.paragraphs[0].font.size = Pt(12)

    def _add_closing_slide(self, presentation, analytics_data):
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        self._add_footer_bar(slide)
        text_box = slide.shapes.add_textbox(Inches(1.1), Inches(2.5), Inches(11), Inches(2))
        frame = text_box.text_frame
        frame.text = "Generated by Questiz"
        frame.paragraphs[0].font.size = Pt(30)
        frame.paragraphs[0].font.bold = True
        frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        detail = frame.add_paragraph()
        detail.text = analytics_data["generated_at"].strftime("%Y-%m-%d %H:%M UTC")
        detail.font.size = Pt(16)
        detail.alignment = PP_ALIGN.CENTER

    def _add_slide_title(self, slide, title):
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.55), Inches(8.5), Inches(0.6))
        frame = title_box.text_frame
        frame.text = title
        frame.paragraphs[0].font.size = Pt(22)
        frame.paragraphs[0].font.bold = True

    def _add_metric_box(self, slide, left, top, width, height, label, value):
        shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = self.brand_color
        shape.line.color.rgb = self.brand_color

        text_frame = shape.text_frame
        text_frame.text = label
        text_frame.paragraphs[0].font.size = Pt(13)
        text_frame.paragraphs[0].font.bold = True
        text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

        paragraph = text_frame.add_paragraph()
        paragraph.text = str(value)
        paragraph.font.size = Pt(24)
        paragraph.font.bold = True
        paragraph.font.color.rgb = RGBColor(255, 255, 255)

    def _add_stats_box(self, slide, section):
        shape = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(8.8),
            Inches(1.5),
            Inches(3.4),
            Inches(4.8),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(248, 250, 252)
        shape.line.color.rgb = self.brand_color
        frame = shape.text_frame
        frame.text = "Key Stats"
        frame.paragraphs[0].font.size = Pt(14)
        frame.paragraphs[0].font.bold = True

        for item in section["highlights"][:6]:
            paragraph = frame.add_paragraph()
            paragraph.text = f"{item['label']}: {item['value']}"
            paragraph.font.size = Pt(12)

    def _add_chart(self, slide, section):
        chart_data = CategoryChartData()
        chart_data.categories = [row["label"] for row in section["table_rows"][:10]]
        chart_data.add_series("Responses", [float(row["count"] or 0) for row in section["table_rows"][:10]])

        chart_type = XL_CHART_TYPE.COLUMN_CLUSTERED
        if section["analytics_type"] == "categorical" and len(section["table_rows"]) <= 6:
            chart_type = XL_CHART_TYPE.PIE

        chart = slide.shapes.add_chart(
            chart_type,
            Inches(0.8),
            Inches(1.6),
            Inches(7.4),
            Inches(4.8),
            chart_data,
        ).chart

        if chart.has_legend:
            chart.legend.include_in_layout = False
        if chart.series:
            chart.series[0].format.fill.solid()
            chart.series[0].format.fill.fore_color.rgb = self.brand_color

    def _add_table_or_text(self, slide, section):
        if section["text_responses"]:
            text_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(7.6), Inches(4.8))
            frame = text_box.text_frame
            frame.text = "Top Responses"
            frame.paragraphs[0].font.size = Pt(14)
            frame.paragraphs[0].font.bold = True
            for response in section["text_responses"][:8]:
                paragraph = frame.add_paragraph()
                paragraph.text = f"• {response}"
                paragraph.font.size = Pt(12)
            return

        rows = min(max(len(section["table_rows"]), 1) + 1, 10)
        cols = max(len(section["table_columns"]), 1)
        table = slide.shapes.add_table(rows, cols, Inches(0.8), Inches(1.6), Inches(7.6), Inches(4.8)).table
        for column_index, title in enumerate(section["table_columns"][:cols]):
            table.cell(0, column_index).text = str(title)

        if section["analytics_type"] in {"matrix", "open_ended", "matrix_plus"}:
            for row_index, row in enumerate(section["table_rows"][: rows - 1], start=1):
                table.cell(row_index, 0).text = row["label"]
                for column_index, value in enumerate(row["values"][: cols - 1], start=1):
                    table.cell(row_index, column_index).text = str(value)
        elif section["analytics_type"] == "demographics":
            current_row = 1
            for field_section in section["field_sections"]:
                if current_row >= rows:
                    break
                table.cell(current_row, 0).text = field_section["field_name"]
                current_row += 1
                for item in field_section["rows"]:
                    if current_row >= rows:
                        break
                    table.cell(current_row, 0).text = item["label"]
                    if cols > 1:
                        table.cell(current_row, 1).text = str(item["count"])
                    if cols > 2:
                        table.cell(current_row, 2).text = str(item["percentage"])
                    current_row += 1
        else:
            for row_index, row in enumerate(section["table_rows"][: rows - 1], start=1):
                table.cell(row_index, 0).text = str(row["label"])
                if cols > 1:
                    table.cell(row_index, 1).text = str(row["count"])
                if cols > 2:
                    table.cell(row_index, 2).text = str(row.get("percentage", ""))

    def _add_footer_bar(self, slide):
        shape = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            Inches(0),
            Inches(7.15),
            Inches(13.333),
            Inches(0.35),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = self.brand_color
        shape.line.fill.background()
